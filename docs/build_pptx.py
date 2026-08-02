#!/usr/bin/env python3
"""Generates .pptx versions of docs/presentation.html for PowerPoint / Google
Slides (Google Slides imports .pptx directly -- no separate format needed).

Two variants, since there's a real fidelity-vs-editability tradeoff:
  - presentation_screenshots.pptx: one full-slide image per slide, rendered
    from the actual HTML via headless Chrome. Pixel-identical to the HTML
    deck; not editable text in PowerPoint (edit the HTML and regenerate).
  - presentation_editable.pptx: rebuilt with native text boxes, tables, and
    shapes, so it's directly editable in PowerPoint/Slides. The custom flow
    diagrams are simplified approximations, not a pixel-perfect match.

Not a runtime dependency of the app -- `pip install python-pptx` into the
venv before running this (not added to requirements.txt).

Usage:
    ./venv/bin/python3 docs/build_pptx.py
"""

import shutil
import subprocess
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

DOCS_DIR = Path(__file__).parent.resolve()
HTML_PATH = DOCS_DIR / "presentation.html"
SLIDE_COUNT = 16
MOCKUP_IMAGE = DOCS_DIR / "assets" / "openclaw_conversation.png"

# ---- palette (matches presentation.html / the dashboard's light theme) ----
BG = RGBColor(0xF8, 0xFA, 0xFC)
PANEL = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)
TEXT = RGBColor(0x1E, 0x29, 0x3B)
TEXT_SECONDARY = RGBColor(0x64, 0x74, 0x8B)
ACCENT = RGBColor(0x25, 0x63, 0xEB)
ACCENT_TINT = RGBColor(0xEF, 0xF6, 0xFF)
SUCCESS = RGBColor(0x05, 0x96, 0x69)
SUCCESS_TINT = RGBColor(0xEC, 0xFD, 0xF5)
WARNING = RGBColor(0xD9, 0x77, 0x06)
WARNING_TINT = RGBColor(0xFF, 0xFB, 0xEB)
DANGER = RGBColor(0xDC, 0x26, 0x26)

FONT = "Inter"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# =========================================================================
# Variant 1: image-based (screenshots of the real HTML deck)
# =========================================================================


def build_image_pptx() -> Path:
    out_path = DOCS_DIR / "presentation_screenshots.pptx"

    # Render the quality-trend chart with its intro animation already
    # settled -- Chart.js animates the line in on load, and a screenshot
    # taken mid-animation would capture the wrong (partial) shape.
    with tempfile.TemporaryDirectory() as tmp:
        qa_html = Path(tmp) / "presentation_qa.html"
        content = HTML_PATH.read_text()
        content = content.replace(
            "responsive: true,\n            maintainAspectRatio: false,",
            "responsive: true,\n            maintainAspectRatio: false,\n            animation: false,",
        )
        qa_html.write_text(content)

        shots = []
        for n in range(1, SLIDE_COUNT + 1):
            shot = Path(tmp) / f"slide_{n:02d}.png"
            subprocess.run(
                [
                    "google-chrome", "--headless=new", "--disable-gpu", "--hide-scrollbars",
                    "--window-size=1920,1080", f"--screenshot={shot}",
                    f"file://{qa_html}#{n}",
                ],
                check=True, capture_output=True,
            )
            shots.append(shot)

        prs = Presentation()
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H
        blank = prs.slide_layouts[6]
        for shot in shots:
            slide = prs.slides.add_slide(blank)
            slide.shapes.add_picture(str(shot), 0, 0, width=SLIDE_W, height=SLIDE_H)
        prs.save(out_path)

    return out_path


# =========================================================================
# Variant 2: native editable slides
# =========================================================================


def _blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    return slide


def _textbox(slide, left, top, width, height, text, *, size=14, bold=False, color=TEXT,
             align=PP_ALIGN.LEFT, font=FONT, italic=False, line_spacing=1.15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.line_spacing = line_spacing
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.name = font
            run.font.color.rgb = color
    return box


def _kicker(slide, text, left=Inches(0.6), top=Inches(0.5)):
    width = Inches(0.35 * len(text) + 0.5)
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.4))
    pill.adjustments[0] = 0.5
    pill.fill.solid()
    pill.fill.fore_color.rgb = ACCENT_TINT
    pill.line.fill.background()
    pill.shadow.inherit = False
    tf = pill.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = Inches(0.12)
    p = tf.paragraphs[0]
    p.text = text.upper()
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.name = FONT
    run.font.color.rgb = ACCENT


def _title(slide, text, top=Inches(1.0)):
    _textbox(slide, Inches(0.6), top, Inches(12.1), Inches(0.9), text, size=30, bold=True, color=TEXT)


def _body(slide, text, top=Inches(1.85), width=Inches(12.1), size=15):
    _textbox(slide, Inches(0.6), top, width, Inches(0.9), text, size=size, color=TEXT_SECONDARY)


def _card(slide, left, top, width, height, heading, body_text, *, accent=None, heading_color=TEXT):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.adjustments[0] = 0.045
    card.fill.solid()
    card.fill.fore_color.rgb = accent[1] if accent else PANEL
    card.line.color.rgb = accent[0] if accent else BORDER
    card.line.width = Pt(1.25)
    card.shadow.inherit = False
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0.18)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = heading
    run = p.runs[0]
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.name = FONT
    run.font.color.rgb = heading_color
    p2 = tf.add_paragraph()
    p2.text = body_text
    p2.space_before = Pt(6)
    p2.line_spacing = 1.2
    for run in p2.runs:
        run.font.size = Pt(12)
        run.font.name = FONT
        run.font.color.rgb = TEXT_SECONDARY
    return card


def _table(slide, left, top, width, height, headers, rows, col_widths=None):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    gshape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = gshape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0xF1, 0xF5, 0xF9)
        p = cell.text_frame.paragraphs[0]
        p.runs[0].font.size = Pt(11)
        p.runs[0].font.bold = True
        p.runs[0].font.name = FONT
        p.runs[0].font.color.rgb = TEXT_SECONDARY
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = PANEL
            p = cell.text_frame.paragraphs[0]
            p.runs[0].font.size = Pt(11.5)
            p.runs[0].font.name = FONT
            p.runs[0].font.color.rgb = TEXT
    return gshape


def _pipeline(slide, top, stages, *, height=Inches(1.15)):
    n = len(stages)
    gap = Inches(0.3)
    total_w = Inches(12.1)
    box_w = Emu(int((total_w - gap * (n - 1)) / n))
    left = Inches(0.6)
    for i, (name, note, kind) in enumerate(stages):
        fill, line = {
            "free": (SUCCESS_TINT, SUCCESS),
            "cheap": (ACCENT_TINT, ACCENT),
            "premium": (WARNING_TINT, WARNING),
        }[kind]
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, box_w, height)
        box.adjustments[0] = 0.08
        box.fill.solid()
        box.fill.fore_color.rgb = fill
        box.line.color.rgb = line
        box.line.width = Pt(1.5)
        box.shadow.inherit = False
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = name
        p.alignment = PP_ALIGN.CENTER
        p.runs[0].font.size = Pt(13)
        p.runs[0].font.bold = True
        p.runs[0].font.name = FONT
        p.runs[0].font.color.rgb = TEXT
        p2 = tf.add_paragraph()
        p2.text = note
        p2.alignment = PP_ALIGN.CENTER
        for run in p2.runs:
            run.font.size = Pt(10)
            run.font.name = FONT
            run.font.color.rgb = TEXT_SECONDARY
        if i < n - 1:
            arrow_left = Emu(int(left + box_w))
            _textbox(slide, arrow_left, Emu(int(top + height / 2 - Inches(0.15))), gap, Inches(0.3),
                     "→", size=16, color=TEXT_SECONDARY, align=PP_ALIGN.CENTER)
        left = Emu(int(left + box_w + gap))


def _footnote(slide, text, top=Inches(6.9)):
    _textbox(slide, Inches(0.6), top, Inches(12.1), Inches(0.5), text, size=10.5,
             color=TEXT_SECONDARY, italic=True)


def build_native_pptx() -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1. Title
    s = _blank_slide(prs)
    logo = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.17), Inches(1.6), Inches(1), Inches(1))
    logo.adjustments[0] = 0.25
    logo.fill.solid(); logo.fill.fore_color.rgb = ACCENT; logo.line.fill.background(); logo.shadow.inherit = False
    _textbox(s, Inches(1.5), Inches(2.9), Inches(10.33), Inches(1.0), "Email Triage Engine",
             size=44, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    _textbox(s, Inches(2.17), Inches(3.85), Inches(9), Inches(1.0),
             "A token-efficient, multi-stage pipeline that filters noise for free, classifies the rest with "
             "a cheap model, and reserves premium LLM calls for the handful of emails that actually matter.",
             size=16, color=TEXT_SECONDARY, align=PP_ALIGN.CENTER)
    _textbox(s, Inches(2.17), Inches(4.7), Inches(9), Inches(0.5),
             "Tiered triage pipeline   ·   Reranker-based routing   ·   Nightly quality audits",
             size=12.5, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    # 2. The problem
    s = _blank_slide(prs)
    _kicker(s, "The problem"); _title(s, "You can't afford to run every email through a premium LLM")
    _body(s, "A naive pipeline sends every unread message's full body to a summarization model. Most inbox "
             "volume is noise — but the cost and latency scale with total volume, not with what's actually "
             "important.")
    _card(s, Inches(0.6), Inches(3.0), Inches(5.85), Inches(1.7), "Naive approach",
          "1,000 unread emails → 1,000 full-body fetches → 1,000 premium LLM calls, every single sync.",
          accent=(DANGER, PANEL), heading_color=DANGER)
    _card(s, Inches(6.85), Inches(3.0), Inches(5.85), Inches(1.7), "Tiered approach",
          "1,000 unread emails → most resolved for free or by a cheap model → only the truly important "
          "ones ever reach the premium summarizer.", accent=(SUCCESS, PANEL), heading_color=SUCCESS)
    _footnote(s, "Every processed message is cached by Message-ID — re-running the pipeline on mail it's "
                 "already seen costs zero additional tokens.")

    # 3. Pipeline overview
    s = _blank_slide(prs)
    _kicker(s, "The solution"); _title(s, "Escalating tiers — cheapest checks run first")
    _body(s, "Every unread email is evaluated in order until a level is decided. Each later stage costs more "
             "tokens and latency, so it only ever runs on whatever the cheaper stages couldn't confidently resolve.")
    _pipeline(s, Inches(2.9), [
        ("Stage 0\nStatic filter", "regex / blacklist — free", "free"),
        ("Stage 0.5\nReranker router", "semantic score — cheap", "cheap"),
        ("Stage 1\nLLM classify", "ternary decision — cheap", "cheap"),
        ("Stage 2\nPremium summary", "full body — expensive", "premium"),
    ])
    _footnote(s, "A VIP-sender bypass and a low-confidence escalation safety net sit alongside this main line.",
              top=Inches(4.4))

    # 4. VIP + Level 0
    s = _blank_slide(prs)
    _kicker(s, "Stage 0"); _title(s, "VIP bypass & the static noise filter")
    _card(s, Inches(0.6), Inches(2.0), Inches(5.85), Inches(1.9), "VIP bypass",
          "If the sender matches whitelist_vip_senders, skip straight to Stage 2 — full body fetch and a "
          "premium summary. No triage decision needed for people who always matter.")
    _card(s, Inches(6.85), Inches(2.0), Inches(5.85), Inches(1.9), "Level 0 — static, free",
          "Regex/substring match against blacklist_keywords / blacklist_senders, unless the sender's domain is "
          "whitelisted. A hit tags the message low / level 0 — no network call, no LLM call.")
    _card(s, Inches(0.6), Inches(4.1), Inches(12.1), Inches(1.1), "Example blacklist keywords",
          "unsubscribe · newsletter · promotions · marketing · no-reply · noreply · "
          "digest · advertisement")

    # 5. Rerank noise filter mechanics
    s = _blank_slide(prs)
    _kicker(s, "Stage 0.5 — deep dive"); _title(s, "How the rerank noise filter actually works")
    _body(s, "Not an LLM call — a semantic reranker (a Cohere/Jina-style /rerank endpoint) scores each "
             "email's relevance against a single fixed “noise” anchor document. One-directional by design: "
             "it only ever skips Level 1 for confidently-noise mail; anything else always still gets a "
             "real Level 1 classification.")
    _card(s, Inches(0.6), Inches(2.75), Inches(12.1), Inches(1.15),
          "“Noise” anchor",
          "“An automated system notification, media download alert, promotional marketing email, "
          "newsletter, or subscription update that does not require any reply or action from you.”",
          accent=(TEXT_SECONDARY, PANEL))
    _table(s, Inches(0.6), Inches(4.15), Inches(12.1), Inches(2.4),
           ["Step", "What happens"],
           [
               ("1. Build the query", "From: {sender} | Subject: {subject} | Snippet: {snippet}"),
               ("2. Call the reranker", 'POST /rerank {model, query, documents:[noise]} '
                                        '→ {results:[{index, relevance_score}]}'),
               ("3. Route on threshold", "noise ≥ 0.999 → Level 0 (skip Level 1, free)  ·  "
                                          "otherwise → Level 1 (normal LLM classification)"),
           ],
           col_widths=[Inches(2.6), Inches(9.5)])

    # 6. Level 1
    s = _blank_slide(prs)
    _kicker(s, "Stage 1"); _title(s, "Cheap LLM ternary classification")
    _body(s, "Sends only From / Subject / Snippet — never the full body — to a cheap model via an "
             "OpenAI-compatible /chat/completions proxy, expecting a strict JSON decision.")
    _table(s, Inches(0.6), Inches(2.7), Inches(12.1), Inches(2.4),
           ["Field", "Meaning"],
           [
               ("suggested_level", "0 (noise) · 1 (notification) · 2 (important)"),
               ("reason", "short natural-language justification"),
               ("confidence_score", "0.0 – 1.0"),
               ("tag", "one-word classification, e.g. promotion, personal, vip"),
           ],
           col_widths=[Inches(3.2), Inches(8.9)])
    _footnote(s, "The same rerank-based classifier from Stage 0.5 can also stand in here (triage_type: \"tei\") "
                 "instead of an LLM call.", top=Inches(5.4))

    # 7. Escalation + Level 2
    s = _blank_slide(prs)
    _kicker(s, "Safety net & Stage 2"); _title(s, "Ambiguity escalation, then the premium summary")
    _card(s, Inches(0.6), Inches(2.0), Inches(5.85), Inches(2.0), "Ambiguity escalation",
          "If Stage 1's confidence_score falls below confidence_threshold (default 0.8), the full body is "
          "fetched and re-evaluated by the premium model — a rare safety net, not the normal path.")
    _card(s, Inches(6.85), Inches(2.0), Inches(5.85), Inches(2.0), "Level 2 — premium summarization",
          "Only for messages that land here: full body → premium model → a bulleted executive "
          "summary, plus its own confidence score and tag.")
    _card(s, Inches(0.6), Inches(4.2), Inches(12.1), Inches(1.0), "SummaryResult schema",
          "summary (string) · confidence_score (0.0–1.0) · tag (e.g. vip)")

    # 8. Caching
    s = _blank_slide(prs)
    _kicker(s, "Cost control"); _title(s, "Cache everything — pay for a message exactly once")
    _body(s, "Every processed message is keyed by its RFC Message-ID in a SQLite email_cache table: triage "
             "level, tag, reason, score, summary, and per-stage token/duration metrics.")
    tiles = [
        ("0", "tokens spent re-processing an already-cached message"),
        ("1", "row per message, updated in place as it moves through stages"),
        ("2", "write phases — download (metadata/body) and triage (decision)"),
    ]
    for i, (num, label) in enumerate(tiles):
        left = Inches(0.6 + i * 4.1)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.9), Inches(3.85), Inches(1.6))
        card.adjustments[0] = 0.06
        card.fill.solid(); card.fill.fore_color.rgb = PANEL
        card.line.color.rgb = BORDER; card.line.width = Pt(1.25); card.shadow.inherit = False
        tf = card.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.text = num; p.alignment = PP_ALIGN.CENTER
        p.runs[0].font.size = Pt(30); p.runs[0].font.bold = True
        p.runs[0].font.name = FONT; p.runs[0].font.color.rgb = ACCENT
        p2 = tf.add_paragraph(); p2.text = label; p2.alignment = PP_ALIGN.CENTER
        for run in p2.runs:
            run.font.size = Pt(11); run.font.name = FONT; run.font.color.rgb = TEXT_SECONDARY
    _footnote(s, "A row can exist before triage runs (metadata downloaded, not yet classified) — "
                 "“processed” specifically means triage_level IS NOT NULL.", top=Inches(4.9))

    # 9. Quality check intro
    s = _blank_slide(prs)
    _kicker(s, "Production quality check")
    _title(s, "How do you know accuracy holds — without reading every email?")
    _body(s, "Cheap models drift silently. A nightly “no-look” audit samples already-triaged mail, "
             "has a separate judge model independently re-decide, and compares the two — no human has to "
             "read anything.")
    _card(s, Inches(0.6), Inches(3.0), Inches(5.85), Inches(2.1), "Stratified sampling",
          "Sampled per triage level (0/1/2), not as one flat draw — so a small sample rate can't miss an "
          "entire level by chance. Pooled across a user's accounts so a low-volume mailbox isn't rounded to zero.")
    _card(s, Inches(6.85), Inches(3.0), Inches(5.85), Inches(2.1), "Example — 10% sample rate",
          "24h window: 20 level-0 + 50 level-1 + 30 level-2 emails → sampled: 2 + 5 + 3 — matching "
          "each level's share, minimum 1.")

    # 10. Quality check flow + metrics
    s = _blank_slide(prs)
    _kicker(s, "Production quality check"); _title(s, "Judge, compare, measure")
    _pipeline(s, Inches(2.0), [
        ("Sample", "stratified by level, pooled across accounts", "cheap"),
        ("Judge model", "independent re-classification", "premium"),
        ("Compare", "judge vs. cached production result", "cheap"),
        ("Metrics", "stored + trended nightly", "free"),
    ], height=Inches(1.0))
    _card(s, Inches(0.6), Inches(3.5), Inches(5.85), Inches(1.9), "Precision / Recall / F1",
          "Macro-averaged over levels {0, 1, 2} — the judge's re-derived level is treated as ground truth, "
          "production's cached level as the prediction being evaluated.")
    _card(s, Inches(6.85), Inches(3.5), Inches(5.85), Inches(1.9), "Summary quality",
          "For messages production actually summarized, the judge grades that summary 1–10 on accuracy, "
          "conciseness, and actionability.")

    # 11. Dashboard (native chart)
    s = _blank_slide(prs)
    _kicker(s, "Admin dashboard"); _title(s, "7-day quality trend, at a glance")
    _body(s, "Runs nightly (default 01:00 UTC), on-demand via Run now, or backfilled historically with a "
             "standalone script. Every run is logged so “nothing ran” is never confused with “it "
             "ran and failed”.")
    chart_data = CategoryChartData()
    chart_data.categories = ["Mon", "Tue", "Wed", "Thu", "Fri", "Sat", "Sun"]
    chart_data.add_series("Level F1 (vs. judge)", (0.78, 0.81, 0.80, 0.83, 0.86, 0.85, 0.88))
    chart_data.add_series("Summary quality (÷10)", (0.70, 0.72, 0.74, 0.73, 0.76, 0.78, 0.79))
    gframe = s.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS, Inches(1.5), Inches(2.7), Inches(10.33), Inches(4.0), chart_data
    )
    chart = gframe.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    plot = chart.plots[0]
    plot.series[0].format.line.color.rgb = ACCENT
    plot.series[1].format.line.color.rgb = WARNING
    chart.value_axis.minimum_scale = 0
    chart.value_axis.maximum_scale = 1

    # 12. Interfaces overview
    s = _blank_slide(prs)
    _kicker(s, "Interfaces"); _title(s, "Three ways to talk to the pipeline")
    _body(s, "The same tiered engine underneath — main.py and mcp_server.py share the same triage "
             "primitives — is exposed three different ways depending on who (or what) is asking.")
    _card(s, Inches(0.6), Inches(2.9), Inches(3.85), Inches(2.9), "CLI",
          "main.py / email-triage.sh — silent JSON-array output by default, built for piping into other "
          "tooling. --human for a rich terminal UI, --level 2 --compact for a minified schema tuned for LLM "
          "agent consumption.")
    _card(s, Inches(4.6), Inches(2.9), Inches(3.85), Inches(2.9), "MCP server",
          "mcp_server.py — Model Context Protocol server for AI agents/editors. stdio transport for local "
          "tools (e.g. Claude Desktop); SSE/HTTP transport for a long-running deployment with the background "
          "scheduler and dashboard.")
    _card(s, Inches(8.6), Inches(2.9), Inches(4.1), Inches(2.9), "Dashboard",
          "A Vite + React SPA — login, per-account sync controls, token usage, integrations, admin "
          "settings, and the quality-check trend — talking to the same Starlette routes the MCP server "
          "exposes.")
    _footnote(s, "MCP access is per-user, token-based (issued at /api/me/mcp-tokens); the dashboard uses "
                 "cookie sessions.", top=Inches(6.1))

    # 13. MCP tool list
    s = _blank_slide(prs)
    _kicker(s, "MCP server"); _title(s, "What an AI agent can actually call")
    _body(s, "Eight tools, split between reading cached state (free) and taking action (mutates a mailbox or "
             "triggers a live sync).", top=Inches(1.8))
    _table(s, Inches(0.6), Inches(2.6), Inches(5.85), Inches(3.6),
           ["Read-only tool", "What it does"],
           [
               ("fetch_and_process_unread", "Cached triage/summaries for unread mail — zero token cost"),
               ("search_emails", "Live Gmail/IMAP search, enriched with cached triage data"),
               ("get_last_download_time", "Per-account sync freshness"),
           ],
           col_widths=[Inches(2.5), Inches(3.35)])
    _table(s, Inches(6.85), Inches(2.6), Inches(5.85), Inches(4.4),
           ["Action tool", "What it does"],
           [
               ("trigger_download", "Forces an immediate sync + triage now"),
               ("mark_emails_as_read", "By triage level, a message, or all"),
               ("create_new_draft", "Creates a new draft (Gmail or IMAP)"),
               ("create_draft_reply", "Drafts a reply to an existing email"),
               ("send_email_reply", "Sends a reply directly, no draft step"),
           ],
           col_widths=[Inches(2.5), Inches(3.35)])

    # 14. MCP client in the wild
    s = _blank_slide(prs)
    _kicker(s, "MCP in the wild"); _title(s, "Any MCP-compatible client can drive this")
    _body(s, "Point a client like OpenClaw at the SSE endpoint with a per-user token (issued from "
             "/api/me/mcp-tokens), and it gets the same 8 tools an editor or custom agent would — including "
             "cached, zero-token-cost reads like fetch_and_process_unread.")
    if MOCKUP_IMAGE.exists():
        pic = s.shapes.add_picture(str(MOCKUP_IMAGE), 0, Inches(2.5), height=Inches(4.1))
        pic.left = Emu(int((SLIDE_W - pic.width) / 2))
        _footnote(s, "Illustrative mockup UI, not an actual OpenClaw screenshot — built to show the shape of "
                     "a real tool-call exchange (request → cached result → synthesized answer).",
                  top=Inches(6.75))
    else:
        _footnote(s, f"[Missing image: {MOCKUP_IMAGE}]", top=Inches(2.6))

    # 15. Architecture recap
    s = _blank_slide(prs)
    _kicker(s, "Under the hood"); _title(s, "Architecture at a glance")
    cells = [
        ("Ingestion & pipeline", "Gmail API + Zoho/IMAP (XOAUTH2), concurrent per account. CLI (main.py) and "
                                  "MCP server (mcp_server.py, FastMCP) share the same tiered pipeline primitives."),
        ("Storage", "SQLite throughout — a per-user email_cache.db, plus a shared app.db for users, "
                     "sessions, integrations, settings, and quality-check history."),
        ("Multi-tenant", "DB-backed accounts with real login, per-user multi-account Gmail/Zoho/IMAP OAuth "
                          "integrations, and per-user MCP tokens for agent access."),
        ("Dashboard", "A Vite + React SPA — sync controls, token usage, and the quality-check trend — "
                       "talking to the same Starlette routes the MCP server exposes."),
    ]
    for i, (heading, body_text) in enumerate(cells):
        left = Inches(0.6 + (i % 2) * 6.25)
        top = Inches(2.0 + (i // 2) * 2.15)
        _card(s, left, top, Inches(5.85), Inches(1.95), heading, body_text)

    # 16. Closing
    s = _blank_slide(prs)
    _textbox(s, Inches(1.17), Inches(3.0), Inches(11), Inches(1.5),
              "Cheap checks first.\nPremium only where it matters.\nQuality measured continuously.",
              size=26, bold=True, color=TEXT, align=PP_ALIGN.CENTER, line_spacing=1.3)
    _textbox(s, Inches(1.17), Inches(4.6), Inches(11), Inches(0.6), "Questions?",
              size=16, color=TEXT_SECONDARY, align=PP_ALIGN.CENTER)

    out_path = DOCS_DIR / "presentation_editable.pptx"
    prs.save(out_path)
    return out_path


if __name__ == "__main__":
    if shutil.which("google-chrome") is None:
        print("google-chrome not found -- skipping the screenshot-based variant.")
    else:
        p1 = build_image_pptx()
        print(f"Wrote {p1}")
    p2 = build_native_pptx()
    print(f"Wrote {p2}")
